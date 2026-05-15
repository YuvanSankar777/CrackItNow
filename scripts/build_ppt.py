"""
Build the formal corporate-style presentation:
"AI Interview Preparation System using RAG and Vector Database"

Style:
    - Navy / white / grey palette
    - Calibri throughout
    - Minimalistic, generous whitespace
    - Section dividers between major groups
    - Subtle fade transitions on every slide

Run:
    backend/venv/bin/python scripts/build_ppt.py
Outputs:
    AI_Interview_Prep.pptx in the repo root.
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

# ── Corporate palette ──────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x25, 0x40)   # primary navy — deep, professional
NAVY_DARK   = RGBColor(0x05, 0x18, 0x2B)
ACCENT_BLUE = RGBColor(0x1F, 0x4E, 0x79)   # secondary, slightly brighter navy
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG     = RGBColor(0xF6, 0xF8, 0xFA)   # very subtle off-white tint
GREY_BORDER = RGBColor(0xD1, 0xD5, 0xDB)
GREY_LIGHT  = RGBColor(0xE5, 0xE7, 0xEB)
GREY_MUTED  = RGBColor(0x6B, 0x72, 0x80)
GREY_DARK   = RGBColor(0x37, 0x41, 0x51)
TEXT        = RGBColor(0x11, 0x18, 0x27)

FONT = 'Calibri'

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

PROJECT_TITLE  = "AI Interview Preparation System"
PROJECT_SUB    = "Adaptive Coding Interviews with a Generative-AI Interviewer & Live Code Sandbox"
INSTITUTION    = "Department of Computer And Communication Engineering · 2026"

TEAM = [
    "Kabilan J",
    "Vikash N",
    "Santhosh Kumar S",
    "Yuvan Sankar NKR",
]


# ── Building blocks ────────────────────────────────────────────────────────

def add_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False
    return slide


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, font=FONT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_card(slide, x, y, w, h, fill=WHITE, border=GREY_BORDER, border_w=0.75):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_w)
    card.shadow.inherit = False
    return card


def add_hline(slide, x, y, w, color=GREY_LIGHT, weight=1.0):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.shadow.inherit = False
    return line


def add_vline(slide, x, y, h, color=GREY_LIGHT, weight=1.0):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(weight), h)
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.shadow.inherit = False
    return line


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT, bullet="▸  ",
                line_spacing=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = bullet + item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_slide_title(slide, eyebrow, title):
    """Standard top-of-slide treatment: thin navy bar, eyebrow label, big title, hairline."""
    # Top accent bar (4pt)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.shadow.inherit = False

    add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.4),
             eyebrow.upper(), size=11, bold=True, color=ACCENT_BLUE)
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.9),
             title, size=30, bold=True, color=NAVY)
    # hairline under title
    add_hline(slide, Inches(0.6), Inches(1.65), Inches(12.1),
              color=GREY_LIGHT, weight=1.0)


def add_footer(slide, num, total):
    add_hline(slide, Inches(0.6), Inches(7.0), Inches(12.1), color=GREY_LIGHT, weight=0.75)
    add_text(slide, Inches(0.6), Inches(7.08), Inches(8), Inches(0.32),
             f"{PROJECT_TITLE}  ·  RAG + Vector DB + LLM",
             size=9, color=GREY_MUTED)
    add_text(slide, Inches(11.5), Inches(7.08), Inches(1.5), Inches(0.32),
             f"{num} / {total}", size=9, color=GREY_MUTED, align=PP_ALIGN.RIGHT)


def add_fade_transition(slide):
    """Inject a subtle fade transition into the slide XML."""
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    }
    transition_xml = (
        f'<p:transition xmlns:p="{nsmap["p"]}" xmlns:p14="{nsmap["p14"]}" '
        f'spd="med" p14:dur="500"><p:fade/></p:transition>'
    )
    transition = etree.fromstring(transition_xml)
    slide.element.append(transition)


def step_circle(slide, x, y, size, label, fill=NAVY, color=WHITE, font_size=14):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    circ.fill.solid(); circ.fill.fore_color.rgb = fill
    circ.line.fill.background(); circ.shadow.inherit = False
    tf = circ.text_frame; tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(label)
    r.font.name = FONT; r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = color


def section_divider(prs, number, title, subtitle=""):
    """A clean section divider: navy band on the left with the number, title in navy."""
    slide = add_blank_slide(prs)
    # full-bleed left band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), SLIDE_H)
    band.line.fill.background(); band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.shadow.inherit = False
    # Section number, big and faint on the band
    add_text(slide, Inches(0.5), Inches(2.3), Inches(3.5), Inches(2.0),
             f"{number:02d}", size=130, bold=True, color=WHITE, font=FONT)
    add_text(slide, Inches(0.6), Inches(4.5), Inches(3.7), Inches(0.4),
             "SECTION", size=11, bold=True, color=WHITE)
    # Title on the right side
    add_text(slide, Inches(5.0), Inches(3.0), Inches(7.8), Inches(1.0),
             title, size=42, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(5.0), Inches(4.1), Inches(7.8), Inches(1.0),
                 subtitle, size=16, color=GREY_MUTED, italic=True)
    # bottom hairline accent
    add_hline(slide, Inches(5.0), Inches(5.4), Inches(2.5), color=NAVY, weight=2)
    return slide


# ── Slides ─────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = add_blank_slide(prs)

    # ── HERO BLOCK (top 3.6") — full-bleed navy with white title ──
    hero_h = Inches(3.6)
    hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, hero_h)
    hero.line.fill.background(); hero.fill.solid(); hero.fill.fore_color.rgb = NAVY
    hero.shadow.inherit = False

    # Lighter navy band on top edge for layered depth
    edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    edge.line.fill.background(); edge.fill.solid(); edge.fill.fore_color.rgb = ACCENT_BLUE
    edge.shadow.inherit = False

    # Decorative cluster — three white dots top-right (geometric accent)
    for i in range(3):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(11.7 + i * 0.35), Inches(0.55),
                                    Inches(0.18), Inches(0.18))
        opacity_fill = WHITE
        d.fill.solid(); d.fill.fore_color.rgb = opacity_fill
        d.line.fill.background(); d.shadow.inherit = False

    # Vertical accent line on left
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(0.7), Inches(0.9),
                                          Inches(0.04), Inches(2.0))
    accent_line.line.fill.background()
    accent_line.fill.solid(); accent_line.fill.fore_color.rgb = WHITE
    accent_line.shadow.inherit = False

    # Eyebrow label
    add_text(slide, Inches(1.0), Inches(0.9), Inches(10), Inches(0.4),
             "PROJECT  ·  AI · MACHINE LEARNING",
             size=11, bold=True, color=WHITE, font=FONT)

    # Big white title (split across two lines for visual rhythm)
    add_text(slide, Inches(1.0), Inches(1.4), Inches(11.5), Inches(1.0),
             "AI Interview Preparation",
             size=48, bold=True, color=WHITE)
    add_text(slide, Inches(1.0), Inches(2.15), Inches(11.5), Inches(0.9),
             "System",
             size=48, bold=True, color=WHITE)

    # Short white rule under title
    add_hline(slide, Inches(1.0), Inches(3.05), Inches(1.5),
              color=WHITE, weight=2.5)

    # ── BODY (white area below hero) ──
    # Subtitle
    add_text(slide, Inches(1.0), Inches(3.85), Inches(11.5), Inches(0.5),
             PROJECT_SUB, size=18, color=NAVY, italic=True)

    # Team Members card
    card_y = Inches(4.7)
    card_h = Inches(1.7)
    add_card(slide, Inches(1.0), card_y, Inches(11.3), card_h,
             fill=GREY_BG, border=GREY_LIGHT, border_w=0.75)
    # Left accent on the card
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(1.0), card_y, Inches(0.1), card_h)
    accent.line.fill.background()
    accent.fill.solid(); accent.fill.fore_color.rgb = NAVY
    accent.shadow.inherit = False

    add_text(slide, Inches(1.3), card_y + Inches(0.15), Inches(11), Inches(0.35),
             "TEAM MEMBERS", size=11, bold=True, color=NAVY)

    # 2x2 team grid
    name_w = Inches(5.4); name_h = Inches(0.45)
    for i, name in enumerate(TEAM):
        col = i % 2
        row = i // 2
        nx = Inches(1.3) + (name_w + Inches(0.3)) * col
        ny = card_y + Inches(0.6) + (name_h + Inches(0.1)) * row
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, nx,
                                     ny + Inches(0.15), Inches(0.1), Inches(0.1))
        sq.fill.solid(); sq.fill.fore_color.rgb = NAVY
        sq.line.fill.background(); sq.shadow.inherit = False
        add_text(slide, nx + Inches(0.25), ny, name_w - Inches(0.25), name_h,
                 name, size=15, bold=True, color=TEXT)

    # Bottom institution band
    add_hline(slide, Inches(1.0), Inches(6.7), Inches(11.3), color=GREY_LIGHT, weight=0.75)
    add_text(slide, Inches(1.0), Inches(6.82), Inches(11.3), Inches(0.4),
             INSTITUTION, size=11, bold=True, color=NAVY)


def slide_introduction(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Introduction", "The case for an intelligent, adaptive interview platform")
    pts = [
        "Technical interviews are the single biggest filter in tech recruitment.",
        "Most candidates rehearse with static question banks — passive, one-way practice.",
        "Existing platforms grade the final answer, not the candidate's reasoning.",
        "Generic question lists ignore the role, level, and company being targeted.",
        "There is no affordable system that simulates a real, interactive interview.",
    ]
    add_bullets(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.2),
                pts, size=18, color=TEXT)
    # Closing strap
    add_card(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.7),
             fill=GREY_BG, border=NAVY, border_w=0.75)
    add_text(slide, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.5),
             "We need an intelligent, adaptive, and conversational interview system.",
             size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Problem Statement", "Five concrete gaps in today's prep tools")
    pairs = [
        ("No personalization",       "Question banks are one-size-fits-all — no tailoring to role or skill level."),
        ("No real-time feedback",    "Candidates discover their gaps only after the real interview, not during practice."),
        ("Static question pool",     "Hard-coded question lists go stale; no dynamic generation."),
        ("Keyword-only search",      "Lookups match exact words, not meaning — relevant questions are missed."),
        ("No interview simulation",  "Reading questions ≠ defending an answer. The conversation loop is missing."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3)) / 2
    cell_h = Inches(1.05)
    for i, (h, body) in enumerate(pairs):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (cell_w + Inches(0.3)) * col
        y = Inches(2.0) + (cell_h + Inches(0.18)) * row
        add_card(slide, x, y, cell_w, cell_h)
        # left navy stripe
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), cell_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background(); bar.shadow.inherit = False
        add_text(slide, x + Inches(0.3), y + Inches(0.15), cell_w - Inches(0.5), Inches(0.4),
                 h, size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), y + Inches(0.55), cell_w - Inches(0.5), Inches(0.5),
                 body, size=12, color=GREY_DARK)


def slide_objectives(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Objectives", "What the system must achieve")
    objs = [
        ("Adaptive interview engine",   "Generate questions that adjust to role, level, and prior answers."),
        ("Semantic retrieval",          "Use embeddings + vector search to fetch contextually similar questions."),
        ("Real-time feedback",          "Score every answer and explain why, not just pass / fail."),
        ("Conversational simulation",   "Mimic a real interviewer: probe, follow up, recalibrate."),
        ("Improved candidate confidence", "Make practice deliberate, measurable, and repeatable."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3) * 2) / 3
    for i, (h, body) in enumerate(objs[:3]):
        x = Inches(0.6) + (cell_w + Inches(0.3)) * i
        add_card(slide, x, Inches(2.0), cell_w, Inches(2.2))
        step_circle(slide, x + Inches(0.25), Inches(2.2), Inches(0.55), str(i + 1), font_size=14)
        add_text(slide, x + Inches(1.0), Inches(2.25), cell_w - Inches(1.2), Inches(0.5),
                 h, size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), Inches(2.95), cell_w - Inches(0.5), Inches(1.2),
                 body, size=12, color=GREY_DARK)
    cell_w2 = (Inches(12.1) - Inches(0.3)) / 2
    for i, (h, body) in enumerate(objs[3:]):
        x = Inches(0.6) + (cell_w2 + Inches(0.3)) * i
        add_card(slide, x, Inches(4.4), cell_w2, Inches(2.2))
        step_circle(slide, x + Inches(0.25), Inches(4.6), Inches(0.55), str(i + 4), font_size=14)
        add_text(slide, x + Inches(1.0), Inches(4.65), cell_w2 - Inches(1.2), Inches(0.5),
                 h, size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), Inches(5.35), cell_w2 - Inches(0.5), Inches(1.2),
                 body, size=12, color=GREY_DARK)


def slide_proposed(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Proposed Solution", "An end-to-end AI interview platform — talk, code, evaluate, adapt")
    bullets = [
        "Use a generative LLM (Google Gemini) to author questions, score answers, and ask probing follow-ups.",
        "Run candidate code in a self-hosted Piston sandbox supporting Python, JavaScript, Java, and C++.",
        "Drive the conversation with the browser's Web Speech API — text-to-speech for the AI, speech-to-text for the user.",
        "Persist every session, question, answer, and evaluation in a Django backend for analytics.",
        "Deliver a resizable IDE-style frontend in React + Monaco so the experience feels like a real interview.",
    ]
    add_bullets(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.4), bullets, size=17)

    # "Why this stack?" highlight strip
    bar_y = Inches(5.6)
    add_card(slide, Inches(0.8), bar_y, Inches(11.7), Inches(1.1),
             fill=GREY_BG, border=NAVY, border_w=0.75)
    add_text(slide, Inches(1.0), bar_y + Inches(0.18), Inches(11.4), Inches(0.4),
             "Why this stack?", size=12, bold=True, color=NAVY)
    add_text(slide, Inches(1.0), bar_y + Inches(0.5), Inches(11.4), Inches(0.6),
             "Free-tier Gemini for intelligence, self-hosted Piston for safe code execution, browser-native voice — zero vendor lock-in, zero per-call cost.",
             size=13, color=TEXT)


def slide_overview(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "System Overview", "High-level interaction in four steps")
    steps = [
        ("Setup the session",          "Candidate picks a target company and difficulty (Easy / Intermediate / Difficult)."),
        ("AI generates the question",  "Gemini produces a coding problem with starter code and 3+ test cases."),
        ("Code, run, submit",          "Candidate writes code in the IDE; Piston executes it against every test case."),
        ("Evaluate, follow up, adapt", "AI scores the answer, asks one probing follow-up, then advances with adjusted difficulty."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3) * 3) / 4
    for i, (h, body) in enumerate(steps):
        x = Inches(0.6) + (cell_w + Inches(0.3)) * i
        add_card(slide, x, Inches(2.2), cell_w, Inches(4.0))
        step_circle(slide, x + (cell_w - Inches(0.7)) / 2, Inches(2.5),
                    Inches(0.7), str(i + 1), font_size=20)
        add_text(slide, x + Inches(0.3), Inches(3.5), cell_w - Inches(0.6), Inches(0.5),
                 h, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), Inches(4.0), cell_w - Inches(0.6), Inches(2.0),
                 body, size=11, color=GREY_DARK, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "System Architecture", "End-to-end pipeline")

    nodes = [
        ("Candidate",      "Browser"),
        ("React Frontend", "Vite · Monaco · Tailwind"),
        ("Django REST API", "Auth · Sessions · Eval"),
        ("Gemini LLM",     "Question · Score · Follow-up"),
        ("Piston Sandbox", "Docker code runner"),
        ("SQLite Store",   "Sessions · Q · A · Evals"),
        ("Result UI",      "Scorecard · Feedback"),
    ]
    n = len(nodes)
    node_w = Inches(1.55); node_h = Inches(1.4)
    gap = Inches(0.18)
    total = node_w * n + gap * (n - 1)
    sx = (SLIDE_W - total) / 2
    y = Inches(3.0)

    for i, (name, sub) in enumerate(nodes):
        x = sx + (node_w + gap) * i
        # node card with thin top accent
        add_card(slide, x, y, node_w, node_h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, node_w, Inches(0.1))
        bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.shadow.inherit = False
        add_text(slide, x + Inches(0.1), y + Inches(0.3), node_w - Inches(0.2), Inches(0.5),
                 name, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), y + Inches(0.85), node_w - Inches(0.2), Inches(0.45),
                 sub, size=9, color=GREY_MUTED, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax1 = x + node_w
            ax2 = ax1 + gap
            ay = y + node_h / 2
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax1, ay, ax2, ay)
            arrow.line.color.rgb = NAVY
            arrow.line.width = Pt(1.5)
            tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, ax2 - Inches(0.08),
                                          ay - Inches(0.05), Inches(0.1), Inches(0.1))
            tri.rotation = 30
            tri.fill.solid(); tri.fill.fore_color.rgb = NAVY
            tri.line.fill.background(); tri.shadow.inherit = False

    add_text(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.5),
             "Every request flows through Django; AI calls go to Gemini, code runs go to Piston — frontend never talks to either directly.",
             size=13, color=GREY_DARK, align=PP_ALIGN.CENTER, italic=True)
    add_text(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.4),
             "Adapt loop — Gemini's evaluation score feeds back into the next question's difficulty and follow-up choice.",
             size=12, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)


def slide_arch_explain(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Architecture Explained", "Responsibilities of each component")
    items = [
        ("React Frontend",   "Vite dev server, Tailwind for styling, Monaco for the IDE, Web Speech API for voice in/out."),
        ("Django REST API",  "Auth, session orchestration, prompt assembly, scoring blend, integrity tracking."),
        ("Gemini LLM",       "Generates the question + test cases, evaluates the candidate's answer, asks one follow-up, writes the final report."),
        ("Piston Sandbox",   "Self-hosted Docker container that compiles and runs candidate code in an isolated environment."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3)) / 2
    cell_h = Inches(2.1)
    for i, (h, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (cell_w + Inches(0.3)) * col
        y = Inches(2.0) + (cell_h + Inches(0.18)) * row
        add_card(slide, x, y, cell_w, cell_h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), cell_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background(); bar.shadow.inherit = False
        add_text(slide, x + Inches(0.35), y + Inches(0.3), cell_w - Inches(0.6), Inches(0.5),
                 h, size=15, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.35), y + Inches(0.85), cell_w - Inches(0.6), Inches(1.2),
                 body, size=12, color=GREY_DARK)


def slide_tech(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Technologies Used", "The stack at a glance")
    cols = [
        ("Frontend",  ["React 19 + Vite", "Tailwind CSS v4", "Monaco Editor (IDE)", "react-resizable-panels", "Web Speech API"]),
        ("Backend",   ["Django 5 + DRF", "JWT auth (SimpleJWT)", "Python 3.11+"]),
        ("AI & Exec", ["Google Gemini 2.5-flash", "Piston (self-hosted Docker)", "gTTS (server fallback)"]),
        ("Storage",   ["SQLite (dev)", "PostgreSQL-ready", "Models: Session · Question · Answer · Evaluation"]),
    ]
    cell_w = (Inches(12.1) - Inches(0.3) * 3) / 4
    for i, (h, items) in enumerate(cols):
        x = Inches(0.6) + (cell_w + Inches(0.3)) * i
        add_card(slide, x, Inches(2.0), cell_w, Inches(4.7))
        # navy header strip
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), cell_w, Inches(0.55))
        head.line.fill.background(); head.fill.solid(); head.fill.fore_color.rgb = NAVY
        head.shadow.inherit = False
        add_text(slide, x + Inches(0.2), Inches(2.08), cell_w - Inches(0.4), Inches(0.4),
                 h, size=14, bold=True, color=WHITE)
        add_bullets(slide, x + Inches(0.25), Inches(2.75), cell_w - Inches(0.5), Inches(3.7),
                    items, size=12, color=TEXT)


def slide_frameworks(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Frameworks & Tools", "What does what")
    rows = [
        ("FastAPI",                  "Async Python web framework — exposes the interview API"),
        ("Sentence Transformers",    "Generates 384-dim embeddings (e.g. all-MiniLM-L6-v2)"),
        ("Cross Encoder (HF)",       "Reranks retrieval candidates for higher precision"),
        ("Vector DB",                "Indexes embeddings + does k-NN similarity search"),
        ("LLM API (OpenAI / local)", "Generates final questions, evaluations, follow-ups"),
        ("Uvicorn",                  "ASGI server hosting the FastAPI app"),
        ("Hugging Face Hub",         "Source for embedding + reranker model weights"),
    ]
    col_w = [Inches(3.4), Inches(8.7)]
    x_tool = Inches(0.6)
    x_desc = x_tool + col_w[0]

    # Header band
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_tool, Inches(2.0),
                                   sum(col_w, Emu(0)), Inches(0.5))
    head.line.fill.background(); head.fill.solid(); head.fill.fore_color.rgb = NAVY
    head.shadow.inherit = False
    add_text(slide, x_tool + Inches(0.2), Inches(2.08), col_w[0], Inches(0.35),
             "Tool", size=12, bold=True, color=WHITE)
    add_text(slide, x_desc + Inches(0.2), Inches(2.08), col_w[1], Inches(0.35),
             "Role", size=12, bold=True, color=WHITE)

    y = Inches(2.6)
    for i, (tool, desc) in enumerate(rows):
        # zebra stripe for readability
        if i % 2 == 0:
            add_card(slide, x_tool, y, sum(col_w, Emu(0)), Inches(0.55),
                     fill=GREY_BG, border=GREY_BG, border_w=0)
        add_text(slide, x_tool + Inches(0.2), y + Inches(0.13), col_w[0], Inches(0.4),
                 tool, size=12, bold=True, color=NAVY)
        add_text(slide, x_desc + Inches(0.2), y + Inches(0.13), col_w[1], Inches(0.4),
                 desc, size=12, color=GREY_DARK)
        y += Inches(0.55)


def slide_dataset(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Dataset Description", "What we feed the system")

    # Left card: composition
    add_card(slide, Inches(0.6), Inches(2.0), Inches(6), Inches(4.7))
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(6), Inches(0.5))
    head.line.fill.background(); head.fill.solid(); head.fill.fore_color.rgb = NAVY
    head.shadow.inherit = False
    add_text(slide, Inches(0.8), Inches(2.08), Inches(5.6), Inches(0.4),
             "Corpus", size=14, bold=True, color=WHITE)
    add_bullets(slide, Inches(0.85), Inches(2.7), Inches(5.5), Inches(3.8), [
        "1,000+ curated interview questions",
        "Categories — Technical, HR, Coding",
        "Sourced from public banks + manually authored",
        "Each entry: question, ideal answer, difficulty, topic",
        "Stored as JSON; embedded once at ingestion time",
    ], size=13)

    # Right card: schema
    add_card(slide, Inches(6.85), Inches(2.0), Inches(5.85), Inches(4.7))
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(2.0), Inches(5.85), Inches(0.5))
    head.line.fill.background(); head.fill.solid(); head.fill.fore_color.rgb = NAVY
    head.shadow.inherit = False
    add_text(slide, Inches(7.05), Inches(2.08), Inches(5.5), Inches(0.4),
             "Per-question schema", size=14, bold=True, color=WHITE)
    schema = [
        ('question',   'str   — the prompt'),
        ('answer',     'str   — reference answer'),
        ('difficulty', 'enum  — easy / medium / hard'),
        ('topic',      'str   — DSA, OS, DBMS, HR ...'),
        ('category',   'enum  — technical / hr / coding'),
    ]
    y = Inches(2.85)
    for k, v in schema:
        add_text(slide, Inches(7.1), y, Inches(1.6), Inches(0.4),
                 k, size=12, bold=True, color=NAVY, font='Consolas')
        add_text(slide, Inches(8.7), y, Inches(3.9), Inches(0.4),
                 v, size=12, color=GREY_DARK, font='Consolas')
        y += Inches(0.55)


def slide_methodology(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Methodology", "From raw text to ranked response")
    steps = [
        ("Data collection",          "Aggregate Q/A from public banks and curated authoring."),
        ("Embedding generation",     "Encode every question with a sentence-transformer."),
        ("Vector indexing",          "Persist embeddings + metadata in a vector DB."),
        ("Query processing",         "Embed the user's input the same way for fair comparison."),
        ("Reranking",                "Cross-encode the top-K results to lift the most relevant."),
        ("Response generation",      "Pass reranked context to the LLM for a final, grounded answer."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3) * 2) / 3
    cell_h = Inches(2.0)
    for i, (h, body) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = Inches(0.6) + (cell_w + Inches(0.3)) * col
        y = Inches(2.0) + (cell_h + Inches(0.25)) * row
        add_card(slide, x, y, cell_w, cell_h)
        step_circle(slide, x + Inches(0.3), y + Inches(0.3),
                    Inches(0.55), str(i + 1), font_size=15)
        add_text(slide, x + Inches(1.0), y + Inches(0.32), cell_w - Inches(1.2), Inches(0.5),
                 h, size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), y + Inches(1.0), cell_w - Inches(0.6), Inches(1.0),
                 body, size=11, color=GREY_DARK)


def slide_workflow(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Working Flow", "What happens on every interview turn")
    steps = [
        "User submits an answer or asks for the next question.",
        "Backend embeds the input with the sentence transformer.",
        "Vector DB returns top-K semantically similar items.",
        "Cross-encoder reranks them by deeper relevance.",
        "LLM receives reranked context and produces the next question or evaluation.",
        "Response is shown; difficulty is adjusted for the next turn.",
    ]
    # numbered list
    y = Inches(2.1)
    for i, s in enumerate(steps):
        step_circle(slide, Inches(0.7), y, Inches(0.45), str(i + 1), font_size=12)
        add_text(slide, Inches(1.4), y + Inches(0.05), Inches(11.0), Inches(0.5),
                 s, size=15, color=TEXT)
        if i < len(steps) - 1:
            add_hline(slide, Inches(1.4), y + Inches(0.6), Inches(11.0), color=GREY_LIGHT, weight=0.5)
        y += Inches(0.7)


def slide_features(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Key Features", "What the user actually experiences")
    features = [
        ("Adaptive interview engine",     "Difficulty calibrated to each candidate."),
        ("Real-time evaluation",          "Score and feedback returned for every answer."),
        ("Semantic question retrieval",   "Finds questions by meaning, not keywords."),
        ("Reranked top-K precision",      "Cross-encoder filters out near-duplicates."),
        ("Multi-category coverage",       "Technical, HR, and coding in one platform."),
        ("RAG-grounded answers",          "Hallucinations minimised by retrieval context."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3) * 2) / 3
    cell_h = Inches(1.95)
    for i, (h, body) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.6) + (cell_w + Inches(0.3)) * col
        y = Inches(2.0) + (cell_h + Inches(0.2)) * row
        add_card(slide, x, y, cell_w, cell_h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), cell_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background(); bar.shadow.inherit = False
        add_text(slide, x + Inches(0.3), y + Inches(0.25), cell_w - Inches(0.5), Inches(0.5),
                 h, size=13, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), y + Inches(0.85), cell_w - Inches(0.5), Inches(1.0),
                 body, size=11, color=GREY_DARK)


def slide_advantages(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Advantages", "Why this beats traditional prep tools")
    items = [
        "Improves interview performance through deliberate, scored practice.",
        "Provides instant evaluation — no waiting for a mentor or a real interview.",
        "Scales to any number of users; the architecture is stateless per request.",
        "Higher answer accuracy thanks to reranking before LLM generation.",
        "Modular: swap the embedder, reranker, or LLM without touching the API surface.",
        "Data-driven: every interaction is logged and replayable for analytics.",
    ]
    add_bullets(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5),
                items, size=16, line_spacing=12)


def slide_limitations(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Limitations", "Honest tradeoffs")
    items = [
        ("Internet dependency",  "Hosted LLMs and vector DBs require a stable connection."),
        ("LLM quality ceiling",  "Output quality is bounded by the underlying model."),
        ("API costs",            "Hosted LLMs (e.g. GPT-4) charge per token — adds up at scale."),
        ("Cold-start latency",   "First query after a long idle period is slower (model load)."),
        ("Bias inheritance",     "Embedder and LLM may carry biases from their training data."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3)) / 2
    cell_h = Inches(1.1)
    for i, (h, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (cell_w + Inches(0.3)) * col
        y = Inches(2.0) + (cell_h + Inches(0.2)) * row
        add_card(slide, x, y, cell_w, cell_h)
        add_text(slide, x + Inches(0.3), y + Inches(0.15), cell_w - Inches(0.6), Inches(0.45),
                 h, size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), y + Inches(0.55), cell_w - Inches(0.6), Inches(0.55),
                 body, size=11, color=GREY_DARK)


def slide_future(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Future Enhancements", "Where this can go next")
    ideas = [
        ("Voice-based interviews",   "Speech-to-text for input, TTS for the AI interviewer."),
        ("Resume-driven questions",  "Upload CV → AI tailors questions to listed skills / projects."),
        ("Advanced scoring",         "Rubric-based grading: communication, depth, correctness."),
        ("Multi-language support",   "Localise the question bank and the LLM prompts."),
        ("Analytics dashboard",      "Track accuracy, weak topics, and progress over time."),
        ("Live coding sandbox",      "Embed an IDE for code questions with sandboxed execution."),
    ]
    cell_w = (Inches(12.1) - Inches(0.3) * 2) / 3
    cell_h = Inches(1.95)
    for i, (h, body) in enumerate(ideas):
        col = i % 3
        row = i // 3
        x = Inches(0.6) + (cell_w + Inches(0.3)) * col
        y = Inches(2.0) + (cell_h + Inches(0.2)) * row
        add_card(slide, x, y, cell_w, cell_h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), cell_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE
        bar.line.fill.background(); bar.shadow.inherit = False
        add_text(slide, x + Inches(0.3), y + Inches(0.25), cell_w - Inches(0.5), Inches(0.5),
                 h, size=13, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), y + Inches(0.85), cell_w - Inches(0.5), Inches(1.0),
                 body, size=11, color=GREY_DARK)


def slide_conclusion(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "Conclusion", "What we built and why it matters")
    pts = [
        "Designed and implemented a RAG-based AI interview preparation platform.",
        "Combined a vector database, cross-encoder reranker, and an LLM into one pipeline.",
        "Delivered semantically accurate retrieval and grounded, real-time evaluation.",
        "Created a modular foundation — easy to extend with voice, analytics, or new models.",
    ]
    add_bullets(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.4),
                pts, size=17, line_spacing=12)

    # closing strap
    add_card(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.0),
             fill=NAVY, border=NAVY, border_w=0)
    add_text(slide, Inches(1.0), Inches(5.85), Inches(11.4), Inches(0.4),
             "OUTCOME", size=10, bold=True, color=WHITE)
    add_text(slide, Inches(1.0), Inches(6.15), Inches(11.4), Inches(0.5),
             "Candidates get personalised, interactive practice that mirrors a real interview.",
             size=14, bold=True, color=WHITE)


def slide_references(prs):
    slide = add_blank_slide(prs)
    add_slide_title(slide, "References", "Key sources we relied on")
    refs = [
        "Lewis et al. — Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks (2020).",
        "Reimers & Gurevych — Sentence-BERT: Sentence Embeddings using Siamese Networks (EMNLP 2019).",
        "Hugging Face — Sentence Transformers documentation (sbert.net).",
        "FastAPI — Official documentation (fastapi.tiangolo.com).",
        "Pinecone / Weaviate / Chroma — Vector database concepts and engineering blogs.",
        "OpenAI / Mistral — Model and API documentation.",
    ]
    add_bullets(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(5),
                refs, size=14, line_spacing=10)


def slide_thanks(prs):
    slide = add_blank_slide(prs)
    # Full navy split: top half navy, bottom half white
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(4.0))
    band.line.fill.background(); band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.shadow.inherit = False

    add_text(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.4),
             "Thank You", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.5),
             "Questions and discussion welcomed.",
             size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # Team Members band on white half
    add_text(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
             "TEAM MEMBERS", size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_hline(slide, Inches(5.5), Inches(4.95), Inches(2.3), color=NAVY, weight=2)
    add_text(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
             "  ·  ".join(TEAM), size=15, color=TEXT, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             INSTITUTION, size=11, color=GREY_MUTED, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Define the deck. Section dividers are interleaved.
    deck = [
        ('content', slide_title),

        ('section', (1, "Project Overview", "Introduction · Problem · Objectives")),
        ('content', slide_introduction),
        ('content', slide_problem),
        ('content', slide_objectives),

        ('section', (2, "Proposed System", "Solution · Overview · Architecture")),
        ('content', slide_proposed),
        ('content', slide_overview),
        ('content', slide_architecture),
        ('content', slide_arch_explain),

        ('section', (3, "Implementation", "Stack · Dataset · Methodology")),
        ('content', slide_tech),
        ('content', slide_frameworks),
        ('content', slide_dataset),
        ('content', slide_methodology),
        ('content', slide_workflow),

        ('section', (4, "Outcomes", "Features · Advantages · Roadmap")),
        ('content', slide_features),
        ('content', slide_advantages),
        ('content', slide_limitations),
        ('content', slide_future),

        ('section', (5, "Closing", "Conclusion · References · Q&A")),
        ('content', slide_conclusion),
        ('content', slide_references),
        ('content', slide_thanks),
    ]

    for kind, payload in deck:
        if kind == 'section':
            num, title, sub = payload
            section_divider(prs, num, title, sub)
        else:
            payload(prs)

    total = len(prs.slides)
    # Footer + transitions on every slide except the cover and the closing
    for i, slide in enumerate(prs.slides):
        add_fade_transition(slide)
        # Skip footer on cover (0) and final thank-you slide
        if i == 0 or i == total - 1:
            continue
        add_footer(slide, i + 1, total)

    out = Path(__file__).resolve().parent.parent / 'AI_Interview_Prep.pptx'
    prs.save(out)
    print(f'Wrote: {out}  ({out.stat().st_size // 1024} KB, {total} slides)')


if __name__ == '__main__':
    build()
